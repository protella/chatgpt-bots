"""
Document Parsing and Content Extraction Handler
This module provides comprehensive document processing capabilities for the chatbot system,
supporting PDFs, Word documents, Excel spreadsheets, and other common file types.
Designed to maintain full document context without truncation and store complete content.
"""
import re
import base64
import asyncio
import subprocess
import zipfile
from concurrent.futures import ThreadPoolExecutor
from io import BytesIO
from typing import Dict, List, Optional, Any
import pdfplumber
import pypdf
from pdf2image import convert_from_bytes
from docx import Document
from pptx import Presentation
import pandas as pd
from config import config
from logger import LoggerMixin

# Hard rule: file processing never touches disk — all extraction operates on
# in-memory bytes/BytesIO only. (Known exceptions OUTSIDE this module, both in the
# scanned-PDF path: pdf2image's poppler backend uses its own internal temp files to
# render pages, and pytesseract likewise hands each rendered page to the tesseract
# binary via its own internal temp files. We never write those ourselves; the bytes
# we hold stay in memory.)

# Per-file extraction timeout. Module constant for now; can move to env config
# once config.py is available to this layer's changes.
EXTRACTION_TIMEOUT_SECONDS = 30
# Office XML formats are ZIP archives; refuse anything that would decompress
# past this (zip-bomb guard).
MAX_OFFICE_DECOMPRESSED_BYTES = 200 * 1024 * 1024
# Dedicated bounded pool so a slow parse can't exhaust the default executor.
# NOTE: a timed-out parse keeps its worker thread until it finishes — the pool
# being bounded caps how many can be stuck at once. Size via DOC_EXTRACTION_WORKERS
# (default 5): OCR can hold a worker ~1-2 min and queue wait counts against the
# read_document timeout, so concurrent scans need parallel workers.
_EXTRACTION_EXECUTOR = ThreadPoolExecutor(max_workers=config.doc_extraction_workers,
                                          thread_name_prefix="doc-extract")
# ---------------------------------------------------------------------------
# File-type routing — ONE central extension→handler map (F49).
#
# This map is the single source of truth for BOTH admission (is_document_file)
# AND dispatch (safe_extract_content). Deriving both from the same table is what
# closes the old hole where an extension was ADMITTED but then latin-1 decoded
# into mojibake because dispatch only knew a handful of filenames.
#
# Every text-extractable type Slack can send belongs here. Anything that would
# only yield confident mojibake (binaries) or leak secrets is in DENIED_* below,
# which is consulted BEFORE any mimetype-positive admission.
# ---------------------------------------------------------------------------

# Text-family extensions — code, config, web, markup, IDL, diffs, subtitles, logs.
# parse_text CANNOT fail (utf-8 → latin-1 → cp1252 → iso-8859-1), so every correctly
# text-shaped file here degrades to readable prose.
_TEXT_EXTENSIONS = (
    # plain / docs / logs
    '.txt', '.md', '.mdx', '.log', '.out', '.err', '.trace',
    # code
    '.py', '.js', '.jsx', '.ts', '.tsx', '.mjs', '.cjs',
    '.go', '.rs', '.java', '.c', '.cpp', '.cc', '.h', '.hpp',
    '.cs', '.rb', '.php', '.sh', '.bash', '.zsh', '.ps1', '.bat',
    '.swift', '.kt', '.scala', '.pl', '.lua', '.r', '.m', '.dart',
    '.ex', '.erl', '.hs', '.clj', '.jl', '.zig', '.asm',
    '.sql',
    # data / config
    '.json', '.jsonl', '.ndjson',
    '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf', '.properties',
    '.hcl', '.tf', '.tfvars', '.editorconfig',
    '.xml', '.html', '.htm',
    # web
    '.css', '.scss', '.sass', '.less', '.vue', '.svelte', '.astro', '.j2',
    # markup
    '.rst', '.adoc', '.org', '.tex', '.bib',
    # IDL / schema
    '.graphql', '.gql', '.proto', '.thrift', '.avsc', '.xsd',
    # diffs / subtitles
    '.diff', '.patch', '.srt', '.vtt', '.ass',
)

# The central extension→handler map.
EXTENSION_HANDLERS = {ext: 'parse_text' for ext in _TEXT_EXTENSIONS}
EXTENSION_HANDLERS.update({
    '.pdf': 'parse_pdf_structured',
    '.docx': 'parse_docx_structured',
    '.pptx': 'parse_pptx_structured',
    '.xlsx': 'parse_excel_adaptive',
    '.xls': 'parse_excel_adaptive',
    '.csv': 'parse_excel_adaptive',
    # tab/pipe-separated: pandas already sniffs the delimiter in _parse_csv_with_pandas
    '.tsv': 'parse_excel_adaptive',
    '.tab': 'parse_excel_adaptive',
    '.psv': 'parse_excel_adaptive',
    '.rtf': 'parse_rtf',
    '.eml': 'parse_email',
    '.ipynb': 'parse_notebook',
})

# Extensionless files with a well-known identity that Slack delivers as text.
# NOT a blanket "no extension → text" rule (that would admit binary blobs); only
# these exact basenames route to parse_text.
KNOWN_FILENAMES = {
    'dockerfile', 'makefile', 'readme', 'license', 'changelog',
    'gemfile', 'procfile', 'jenkinsfile', '.gitignore', '.bashrc',
}

# Denylist — consulted BEFORE any mimetype-positive admission. A .zip / .env
# mislabeled text/plain must still be refused: latin-1 decodes ANY bytes into
# confident mojibake (a regression, not coverage), and ingesting secrets into a
# Slack thread + document ledger is not something we want to do quietly.
DENIED_EXTENSIONS = (
    # secrets
    '.env', '.pem', '.key',
    # archives
    '.zip', '.rar', '.7z', '.tar', '.gz',
    # legacy binary Office (no no-disk extractor; must WIN over a lying text/* mimetype so
    # they aren't latin-1'd into mojibake). Only the OOXML variants (.docx/.pptx/.xlsx) parse.
    '.doc', '.ppt', '.xlsb',
    # deliberately-dropped container/ebook/message/columnar formats (licensing/memory —
    # see SPEC_ADDENDUM). Denied outright so a mislabeled text/* can't smuggle them in.
    '.odt', '.ods', '.odp', '.epub', '.msg', '.parquet',
    # raster / vector images: these belong to the vision path, never the document parser —
    # a .jpg claiming text/plain would otherwise be latin-1'd into garbage.
    '.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.tif', '.tiff',
    '.heic', '.heif', '.svg', '.ico',
    # binary design / media assets
    '.psd', '.ai', '.sketch', '.fig',
    '.mp4', '.mov', '.avi', '.mkv',
    '.mp3', '.wav', '.m4a', '.flac',
    # executables / disk images / fonts with no text to extract
    '.exe', '.dll', '.so', '.bin', '.iso',
    '.ttf', '.otf', '.woff', '.woff2',
)

# Document file extensions — derived from the central map so admission and
# dispatch can never disagree. Kept as a plain `set` (public: tests import it).
DOCUMENT_EXTENSIONS = set(EXTENSION_HANDLERS)

# Supported document MIME types. A mimetype here is admitted directly (the
# text/* family is additionally admitted via a catch-all in is_document_file).
# Kept a strict subset of MIME_TYPE_HANDLERS — enforced by a parity test.
SUPPORTED_DOCUMENT_MIMETYPES = {
    # PDF documents
    "application/pdf",
    # Word documents (.doc legacy binary removed: its only extractor, docx2txt,
    # requires a real file path, which the no-disk rule forbids)
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # .docx
    # Excel/Spreadsheet documents
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # .xlsx
    "application/vnd.ms-excel",  # .xls
    "text/csv",  # .csv
    "text/tab-separated-values",  # .tsv
    "text/tsv",  # .tsv (alt label)
    # PowerPoint documents (.ppt legacy binary is not parseable by python-pptx
    # and is deliberately unsupported so users get the proper warning)
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
    # Text documents
    "text/plain",  # .txt
    "text/markdown",  # .md
    "application/rtf",  # .rtf
    "text/rtf",  # .rtf (alt label)
    # Rich messages / notebooks
    "message/rfc822",  # .eml
    "application/x-ipynb+json",  # .ipynb
    # Code files
    "text/x-python",  # .py
    "application/javascript",  # .js
    "application/json",  # .json
    "text/x-sql",  # .sql
    "application/x-yaml",  # .yaml/.yml
    "application/xml",  # .xml
    "text/html",  # .html
    # Log files
    "text/x-log",  # .log
}
# MIME type routing handlers — declared-mimetype fallback for extensionless files.
# (Dispatch is extension-first; see safe_extract_content.)
MIME_TYPE_HANDLERS = {
    'application/pdf': 'parse_pdf_structured',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'parse_docx_structured',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'parse_pptx_structured',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'parse_excel_adaptive',
    'application/vnd.ms-excel': 'parse_excel_adaptive',
    'text/csv': 'parse_excel_adaptive',
    'text/tab-separated-values': 'parse_excel_adaptive',
    'text/tsv': 'parse_excel_adaptive',
    'text/plain': 'parse_text',
    'text/markdown': 'parse_text',
    'application/rtf': 'parse_rtf',
    'text/rtf': 'parse_rtf',
    'message/rfc822': 'parse_email',
    'application/x-ipynb+json': 'parse_notebook',
    'text/x-python': 'parse_text',
    'application/javascript': 'parse_text',
    'application/json': 'parse_text',
    'text/x-sql': 'parse_text',
    'application/x-yaml': 'parse_text',
    'application/xml': 'parse_text',
    'text/html': 'parse_text',
    'text/x-log': 'parse_text',
}
class DocumentHandler(LoggerMixin):
    """Handles document parsing and content extraction with error recovery"""
    def __init__(self, max_document_size: int = 50 * 1024 * 1024):
        """
        Initialize the document handler
        Args:
            max_document_size: Maximum document size in bytes (default 50MB)
        """
        self.max_document_size = max_document_size
    # Async entry point (single implementation; per-parser async wrappers were
    # phantoms referencing methods that never existed and have been removed)
    async def safe_extract_content_async(self, file_data: bytes, mime_type: str, filename: str,
                                          ocr_images: bool = True,
                                          ocr_text: bool = False) -> Dict[str, Any]:
        """Async entry point for document extraction.

        Deliberately a thin executor wrapper around the SYNC safe_extract_content:
        a previous hand-rolled async router drifted from the sync routing table
        (missing text types, a phantom CSV handler, no sanitization, no metadata,
        no fallback chain). One implementation, one behavior — the parity unit
        test enforces the routing table itself. The executor offload matters doubly
        for ocr_text=True: OCR is subprocess+CPU heavy and must not block the loop.
        """
        # OCR (subprocess + CPU per page) blows past the 30s non-OCR cap on real scans,
        # so widen the inner executor timeout to match the worst case the outer per-tool
        # timeout now permits — otherwise this cap would abort an OCR run the tool allows.
        extraction_timeout = EXTRACTION_TIMEOUT_SECONDS
        if ocr_text:
            extraction_timeout = max(EXTRACTION_TIMEOUT_SECONDS,
                                     30 + config.ocr_max_pages * 5)
        loop = asyncio.get_running_loop()
        try:
            return await asyncio.wait_for(
                loop.run_in_executor(
                    _EXTRACTION_EXECUTOR,
                    lambda: self.safe_extract_content(file_data, mime_type, filename,
                                                      ocr_images=ocr_images,
                                                      ocr_text=ocr_text),
                ),
                timeout=extraction_timeout,
            )
        except asyncio.TimeoutError:
            self.log_error(f"Extraction timed out after {extraction_timeout}s for {filename}")
            return {
                'content': f'[Unable to parse {filename} - extraction timed out]',
                'filename': filename,
                'mime_type': mime_type,
                'size_bytes': len(file_data),
                'error': f'Extraction timed out after {extraction_timeout}s',
                'format': 'error',
            }
    def is_document_file(self, filename: str, mimetype: Optional[str] = None) -> bool:
        """
        Check if a file is a supported document type.

        Gate order is load-bearing (F49):
          1. Denylist wins over EVERYTHING — a .zip/.env mislabeled text/plain is
             still refused, because latin-1 would decode any bytes into confident
             mojibake and secrets should not be ingested silently.
          2. Mimetype-positive: the explicit supported set OR the text/* family
             catch-all (fixes every correctly-labeled text/x-* code type).
          3. Extension / known-filename. This is where TypeScript survives: Slack
             sends `video/mp2t` for a .ts file, so it fails step 2 and is admitted
             ONLY by its extension. Never refactor step 2 to short-circuit False.

        Args:
            filename: The filename to check (Slack documents this as nullable)
            mimetype: Optional MIME type
        Returns:
            True if the file is a supported document type
        """
        filename_lower = (filename or "").lower()
        if self._is_denied_file(filename_lower):
            return False
        if mimetype and (mimetype in SUPPORTED_DOCUMENT_MIMETYPES
                         or mimetype.startswith("text/")):
            return True
        return self._handler_for_filename(filename_lower) is not None

    @staticmethod
    def _is_denied_file(filename_lower: str) -> bool:
        """True for secrets and binary formats we refuse regardless of mimetype."""
        return any(filename_lower.endswith(ext) for ext in DENIED_EXTENSIONS)

    @staticmethod
    def _handler_for_filename(filename: Optional[str]) -> Optional[str]:
        """Resolve a parser method name from the central extension map.

        Returns the handler for the LONGEST matching extension (so compound names
        like ``report.pdf.txt`` route by their real trailing extension), then falls
        back to a known extensionless basename, then None (caller consults the
        declared mimetype). Kept static + pure so admission and dispatch share it.
        """
        name = (filename or "").lower().rsplit('/', 1)[-1]
        if name in KNOWN_FILENAMES:
            return 'parse_text'
        best_ext: Optional[str] = None
        for ext in EXTENSION_HANDLERS:
            if name.endswith(ext) and (best_ext is None or len(ext) > len(best_ext)):
                best_ext = ext
        return EXTENSION_HANDLERS[best_ext] if best_ext else None
    def safe_extract_content(self, file_data: bytes, mime_type: str, filename: str,
                             ocr_images: bool = True, ocr_text: bool = False) -> Dict[str, Any]:
        """
        Safely extract document content with comprehensive error recovery
        Args:
            file_data: The document file data as bytes
            mime_type: The MIME type of the document
            filename: The original filename
        Returns:
            Dict with extracted content, structure info, and any errors
        """
        # Validate file size
        if len(file_data) > self.max_document_size:
            return {
                'content': f'[Document {filename} too large: {len(file_data) / 1024 / 1024:.1f}MB (max: {self.max_document_size / 1024 / 1024:.1f}MB)]',
                'error': 'Document exceeds size limit',
                'format': 'error'
            }
        try:
            # Route to appropriate parser. Extension-first, then the declared
            # mimetype for extensionless files, then plain text. Extension-first is
            # deliberate: Slack routinely sends application/octet-stream or a generic
            # text/plain, and (for example) an .ipynb mislabeled application/json must
            # route to parse_notebook — NOT dump raw JSON (a context bomb) at the model.
            # The central EXTENSION_HANDLERS map is the same one admission uses, so an
            # admitted extension can never resolve to a different handler here.
            handler_name = self._handler_for_filename(filename)
            if handler_name is None:
                handler_name = MIME_TYPE_HANDLERS.get(mime_type, 'parse_text')
            # Zip-bomb guard for office-XML formats (they are ZIP archives)
            if handler_name in ('parse_docx_structured', 'parse_pptx_structured', 'parse_excel_adaptive') \
                    and file_data[:2] == b'PK' and not self._office_zip_within_limits(file_data):
                self.log_error(f"Refusing {filename}: decompressed size exceeds "
                               f"{MAX_OFFICE_DECOMPRESSED_BYTES // (1024 * 1024)}MB (zip-bomb guard)")
                return {
                    'content': f'[Unable to parse {filename} - archive decompresses beyond safe limits]',
                    'filename': filename,
                    'mime_type': mime_type,
                    'size_bytes': len(file_data),
                    'error': 'Decompressed size exceeds safety limit',
                    'format': 'error',
                }
            # Get the parser method (PDF takes the OCR-images toggle; the native
            # input_file route sets ocr_images=False so scans skip pdf2image entirely)
            if handler_name == 'parse_pdf_structured':
                result = self.parse_pdf_structured(file_data, filename,
                                                   ocr_images=ocr_images, ocr_text=ocr_text)
            else:
                parser_method = getattr(self, handler_name)
                result = parser_method(file_data, filename)
            # Sanitize the content
            if 'content' in result:
                result['content'] = self.sanitize_content(result['content'])
            # Add filename to result
            result['filename'] = filename
            result['mime_type'] = mime_type
            result['size_bytes'] = len(file_data)
            return result
        except Exception as e:
            self.log_error(f"Failed to parse {filename}: {e}", exc_info=True)
            # Fallback: try basic text extraction
            try:
                raw_text = self.force_text_extraction(file_data, mime_type, filename)
                return {
                    'content': self.sanitize_content(raw_text),
                    'filename': filename,
                    'mime_type': mime_type,
                    'size_bytes': len(file_data),
                    'error': f'Partial extraction - document may be malformed: {str(e)}',
                    'format': 'text'
                }
            except Exception as fallback_error:
                self.log_error(f"Fallback extraction also failed for {filename}: {fallback_error}")
                # Last resort: mark as unparseable
                return {
                    'content': f'[Unable to parse {filename} - document appears malformed or corrupted]',
                    'filename': filename,
                    'mime_type': mime_type,
                    'size_bytes': len(file_data),
                    'error': f'Document could not be parsed: {str(e)}',
                    'format': 'error'
                }
    def sanitize_content(self, text: str) -> str:
        """
        Sanitize content to prevent format breaking and injection attacks
        Args:
            text: The text content to sanitize
        Returns:
            Sanitized text content
        """
        if not text:
            return ""
        sanitized = str(text)
        # Remove null bytes and control characters (except newlines, tabs, carriage returns)
        sanitized = sanitized.replace('\x00', '')
        sanitized = ''.join(char for char in sanitized if ord(char) >= 32 or char in '\n\r\t')
        # Balance code blocks
        code_blocks = sanitized.count('```')
        if code_blocks % 2 != 0:
            sanitized += '\n```'  # Close unclosed code block
        # Escape document markers to prevent injection
        sanitized = sanitized.replace('[Document:', '[Document\\:')
        sanitized = sanitized.replace('[End Document]', '[End\\ Document]')
        sanitized = sanitized.replace('[Page ', '[Page\\ ')
        sanitized = sanitized.replace('[Sheet:', '[Sheet\\:')
        # Fix markdown tables
        sanitized = self.fix_markdown_tables(sanitized)
        # Limit consecutive newlines
        sanitized = re.sub(r'\n{4,}', '\n\n\n', sanitized)
        # No size limit - let the model's token limit handle it
        # Previously limited to 1MB
        return sanitized
    def parse_pdf_structured(self, file_data: bytes, filename: str,
                             ocr_images: bool = True, ocr_text: bool = False) -> Dict[str, Any]:
        """
        Extract PDF content preserving tables and page structure
        Args:
            file_data: PDF file data as bytes
            filename: Original filename for error reporting
            ocr_images: For image-based PDFs, render pages to base64 images for the
                vision path (attach-turn big-file local route). Native input sets this
                False so scans skip pdf2image entirely.
            ocr_text: For image-based PDFs, OCR the rendered pages to plain TEXT and
                fold it into result['content'] (gated on config.enable_pdf_ocr). This is
                orthogonal to ocr_images: the read_document text tool wants text only,
                the local attach path wants both text AND page images.
        Returns:
            Dict with pages, content, and structure info
        """
        # Try pdfplumber first, then pypdf as fallback
        try:
            result = self._parse_pdf_with_pdfplumber(file_data, filename)
        except Exception as e:
            self.log_warning(f"pdfplumber failed, trying pypdf: {e}")
            result = self._parse_pdf_with_pypdf2(file_data, filename)
        # Check if PDF is likely image-based (scanned document)
        if result and self._is_image_based_pdf(result):
            result['is_image_based'] = True
            result['requires_ocr'] = True
            # OCR text is orthogonal to the vision page-image path and gated on config.
            do_ocr = ocr_text and config.enable_pdf_ocr
            ocr_pages = (
                self.ocr_pdf_pages(file_data, max_pages=config.ocr_max_pages, dpi=config.ocr_dpi)
                if do_ocr else []
            )
            if not ocr_images:
                # Text-only route (native attach turn / read_document tool): never render
                # page images. If OCR produced text, use it; otherwise fall back to the
                # honest scanned-document note.
                if ocr_pages:
                    ocr_content, ocr_page_entries = self._format_ocr_pages(
                        ocr_pages, result.get('total_pages', 0))
                    result['content'] = ocr_content
                    result['pages'] = ocr_page_entries
                    result['ocr_text_used'] = True
                    return result
                # Caller will present the PDF natively (rendered pages) — skip the
                # pdf2image conversion (and its poppler temp files) entirely.
                result['content'] = (
                    f"[Note: This PDF appears to be a scanned document; text extraction found "
                    f"minimal content. The document is being provided to the model as rendered pages.]\n\n"
                    f"{result.get('content', '[No text extracted from PDF]')}"
                )
                return result
            # Convert PDF pages to images for vision processing
            self.log_info(f"Converting image-based PDF {filename} to images for OCR")
            page_images = self.convert_pdf_to_images(file_data, max_pages=10)
            if page_images:
                result['page_images'] = page_images
                result['ocr_available'] = True
                self.log_info(f"Successfully converted {len(page_images)} PDF pages to images")
                # Update content to indicate OCR is available
                total_pages = result.get('total_pages', 0)
                if total_pages > len(page_images):
                    result['content'] = (
                        f"[Note: This PDF appears to be a scanned document with {total_pages} total pages. "
                        f"Due to API limits, converted first {len(page_images)} pages to images for vision/OCR analysis. "
                        f"Remaining {total_pages - len(page_images)} pages were not processed.]\n\n"
                        f"{result.get('content', '[No text extracted from PDF]')}"
                    )
                else:
                    result['content'] = (
                        f"[Note: This PDF appears to be a scanned document. "
                        f"Converted all {len(page_images)} page(s) to images for vision analysis. "
                        f"The bot will use vision/OCR to extract content.]\n\n"
                        f"{result.get('content', '[No text extracted from PDF]')}"
                    )
            else:
                # Conversion failed, use original note
                result['content'] = (
                    f"[Note: This PDF appears to be a scanned document or contains primarily images. "
                    f"Text extraction found minimal or no text content. "
                    f"PDF to image conversion failed - OCR not available.]\n\n"
                    f"{result.get('content', '[No text extracted from PDF]')}"
                )
            # OCR text ALSO improves summaries for scans on the big-file local path — emit
            # BOTH the page images (above) and the OCR text content when OCR succeeded.
            if ocr_pages:
                ocr_content, ocr_page_entries = self._format_ocr_pages(
                    ocr_pages, result.get('total_pages', 0))
                result['content'] = f"{result['content']}\n\n{ocr_content}"
                result['pages'] = ocr_page_entries
                result['ocr_text_used'] = True
        return result

    def _format_ocr_pages(self, ocr_pages: List[str],
                          total_pages: int) -> tuple:
        """Build page-structured OCR content + per-page entries from OCR'd page texts.

        Returns (content_str, pages_list). Truncation is LOUD: if the document has more
        pages than were OCR'd (OCR page cap), a bracketed note is prepended saying how
        many of how many pages were read.
        """
        blocks: List[str] = []
        if total_pages and total_pages > len(ocr_pages):
            blocks.append(
                f"[OCR text extracted from the first {len(ocr_pages)} of {total_pages} page(s); "
                f"the remaining {total_pages - len(ocr_pages)} page(s) exceed the OCR page limit "
                f"and were not read.]"
            )
        pages: List[Dict[str, Any]] = []
        for i, page_text in enumerate(ocr_pages, start=1):
            text = (page_text or "").strip() or "[Page contained no OCR-readable text]"
            blocks.append(f"[Page {i}]\n{text}")
            pages.append({'page': i, 'content': text, 'ocr': True})
        return '\n\n'.join(blocks), pages

    def ocr_pdf_pages(self, file_data: bytes, max_pages: int = 20,
                      dpi: int = 300) -> List[str]:
        """OCR an image-only/scanned PDF to per-page plain text.

        Returns one string per rendered page (index 0 == page 1), capped at max_pages.
        Everything stays in memory except the poppler/tesseract temp files documented at
        the module top. NEVER raises: a missing pytesseract package, a missing tesseract
        binary, or any render/OCR error logs a warning and yields [] so the caller falls
        back to the honest scanned-document note.
        """
        try:
            import pytesseract
        except ImportError as e:
            self.log_warning(f"pytesseract not installed; skipping PDF OCR: {e}")
            return []
        try:
            images = convert_from_bytes(file_data, dpi=dpi, fmt='png',
                                        first_page=1, last_page=max_pages)
        except Exception as e:
            self.log_error(f"Failed to render PDF for OCR: {e}")
            if "poppler" in str(e).lower():
                self.log_error("poppler-utils may not be installed. Install with: "
                               "apt-get install poppler-utils (Linux) or brew install poppler (Mac)")
            return []
        page_texts: List[str] = []
        for i, image in enumerate(images[:max_pages], start=1):
            try:
                page_texts.append(pytesseract.image_to_string(image) or "")
            except pytesseract.TesseractNotFoundError as e:
                self.log_warning(f"tesseract binary not found; skipping PDF OCR: {e}")
                return []
            except Exception as e:
                self.log_warning(f"OCR failed on page {i}: {e}")
                page_texts.append("")
        if page_texts:
            self.log_info(f"OCR'd {len(page_texts)} PDF page(s) at {dpi} DPI")
        return page_texts
    def _parse_pdf_with_pdfplumber(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Parse PDF using pdfplumber for advanced structure extraction"""
        pages = []
        total_pages = 0
        has_tables = False
        try:
            with pdfplumber.open(BytesIO(file_data)) as pdf:
                total_pages = len(pdf.pages)
                # Sanity limit for very large PDFs
                pages_to_process = min(total_pages, 1000)
                for i, page in enumerate(pdf.pages[:pages_to_process]):
                    page_data = {'page': i + 1}
                    try:
                        # Extract tables if present
                        tables = page.extract_tables()
                        if tables:
                            page_data['tables'] = []
                            for table in tables:
                                if table and len(table) > 0:
                                    md_table = self.flexible_table_to_markdown(table)
                                    if md_table:
                                        page_data['tables'].append(md_table)
                                        has_tables = True
                        # Extract text content
                        text = page.extract_text() or ""
                        page_data['content'] = text
                        # Try to detect structure hints
                        page_data['structure_hints'] = self._detect_text_structure(text)
                    except Exception as e:
                        self.log_warning(f"Error processing page {i+1} of {filename}: {e}")
                        page_data['content'] = f'[Page {i+1} extraction failed]'
                    pages.append(page_data)
                if pages_to_process < total_pages:
                    pages.append({
                        'page': pages_to_process + 1,
                        'content': f'[Remaining {total_pages - pages_to_process} pages omitted - document too large]'
                    })
        except Exception as e:
            raise Exception(f"pdfplumber parsing failed: {e}")
        # Combine all content
        all_content = []
        for page in pages:
            all_content.append(f"[Page {page['page']}]")
            if page.get('tables'):
                all_content.extend(page['tables'])
            if page.get('content'):
                all_content.append(page['content'])
        return {
            'content': '\n'.join(all_content),
            'pages': pages,
            'total_pages': total_pages,
            'has_tables': has_tables,
            'format': 'pdf'
        }
    def _parse_pdf_with_pypdf2(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Parse PDF using pypdf as fallback"""
        try:
            reader = pypdf.PdfReader(BytesIO(file_data))
            total_pages = len(reader.pages)
            pages = []
            content_parts = []
            # Limit pages for very large documents
            pages_to_process = min(total_pages, 100)
            for i, page in enumerate(reader.pages[:pages_to_process]):
                try:
                    text = page.extract_text()
                    page_data = {
                        'page': i + 1,
                        'content': text or f'[Page {i+1} - no text extracted]'
                    }
                    pages.append(page_data)
                    content_parts.append(f"[Page {i+1}]")
                    content_parts.append(page_data['content'])
                except Exception as e:
                    self.log_warning(f"Error extracting page {i+1}: {e}")
                    pages.append({
                        'page': i + 1,
                        'content': f'[Page {i+1} extraction failed]'
                    })
            if pages_to_process < total_pages:
                content_parts.append(f'[Remaining {total_pages - pages_to_process} pages omitted]')
            return {
                'content': '\n'.join(content_parts),
                'pages': pages,
                'total_pages': total_pages,
                'has_tables': False,
                'format': 'pdf',
                'extraction_method': 'pypdf_fallback'
            }
        except Exception as e:
            raise Exception(f"pypdf parsing failed: {e}")
    def _office_zip_within_limits(self, file_data: bytes) -> bool:
        """Probe an office-XML ZIP's declared decompressed size without extracting."""
        try:
            with zipfile.ZipFile(BytesIO(file_data)) as zf:
                total = sum(info.file_size for info in zf.infolist())
            return total <= MAX_OFFICE_DECOMPRESSED_BYTES
        except zipfile.BadZipFile:
            # Not actually a ZIP despite the PK prefix — let the parser's own
            # error recovery handle it.
            return True

    def parse_pptx_structured(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """
        Extract PowerPoint (.pptx) content: slide text, tables, and speaker notes.
        BytesIO-fed only — no disk access.
        """
        prs = Presentation(BytesIO(file_data))
        content_blocks = []
        tables_found = False
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_blocks = []
            for shape in slide.shapes:
                if getattr(shape, 'has_text_frame', False) and shape.text_frame:
                    text = '\n'.join(
                        run.text for para in shape.text_frame.paragraphs
                        for run in para.runs if run.text
                    ).strip()
                    if not text:
                        # Some decks put text directly on paragraphs without runs
                        text = shape.text_frame.text.strip()
                    if text:
                        slide_blocks.append(text)
                if getattr(shape, 'has_table', False) and shape.has_table:
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    table_md = self.flexible_table_to_markdown(rows)
                    if table_md:
                        slide_blocks.append(table_md)
                        tables_found = True
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_blocks.append(f"Speaker notes: {notes}")
            if slide_blocks:
                content_blocks.append(f"## Slide {slide_num}\n" + '\n\n'.join(slide_blocks))
        full_content = '\n\n'.join(content_blocks)
        return {
            'content': full_content or '[No text content found in presentation]',
            'format': 'pptx',
            'total_slides': len(prs.slides),
            'has_tables': tables_found,
        }

    def parse_docx_structured(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """
        Extract Word document content preserving formatting and structure
        Args:
            file_data: DOCX file data as bytes
            filename: Original filename for error reporting
        Returns:
            Dict with content, pages/sections, and structure info
        """
        try:
            doc = Document(BytesIO(file_data))
            content_blocks = []
            tables_found = False
            # Extract paragraphs and tables in document order
            for element in doc.element.body:
                if element.tag.endswith('p'):  # Paragraph
                    # Find the corresponding paragraph object
                    for para in doc.paragraphs:
                        if para._element == element:
                            text = para.text.strip()
                            if text:  # Only include non-empty paragraphs
                                # Detect if this looks like a header
                                style_name = para.style.name if para.style else ''
                                if 'Heading' in style_name or text.isupper() and len(text) < 100:
                                    content_blocks.append(f"## {text}")
                                else:
                                    content_blocks.append(text)
                            break
                elif element.tag.endswith('tbl'):  # Table
                    # Find the corresponding table object
                    for table in doc.tables:
                        if table._element == element:
                            table_md = self._extract_docx_table(table)
                            if table_md:
                                content_blocks.append(table_md)
                                tables_found = True
                            break
            # Combine content
            full_content = '\n\n'.join(content_blocks)
            # Split into sections if headers are detected
            sections = self._split_into_sections(full_content)
            return {
                'content': full_content,
                'sections': sections,
                'has_tables': tables_found,
                'total_sections': len(sections),
                'format': 'docx'
            }
        except Exception as e:
            # Try alternative extraction method for problematic DOCX files
            self.log_warning(f"Standard DOCX parsing failed, trying alternative method: {e}")
            return self.parse_docx_alternative(file_data, filename)
    def parse_docx_alternative(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """
        Alternative DOCX parsing using ZIP extraction for problematic files
        """
        import zipfile
        import xml.etree.ElementTree as ET
        try:
            # First check if this is actually a ZIP file
            if not file_data.startswith(b'PK'):  # ZIP files start with 'PK'
                self.log_error(f"File {filename} does not appear to be a ZIP/DOCX file. First bytes: {file_data[:4]}")
                # Likely a legacy .doc (OLE2 binary). Unsupported: its only
                # extractor (docx2txt) requires a temp file on disk, which the
                # no-disk rule forbids. Users should re-save as .docx.
                return {
                    'content': f'[{filename} appears to be a legacy .doc file, which is not supported - please re-save it as .docx]',
                    'format': 'doc',
                    'error': 'Legacy .doc format not supported',
                }
            # DOCX files are ZIP archives
            with zipfile.ZipFile(BytesIO(file_data)) as zip_file:
                # Debug: log what's in the ZIP file
                file_list = zip_file.namelist()
                self.log_info(f"DOCX ZIP contents for {filename}: {file_list[:10]}...")  # First 10 files
                # Check if it's a valid DOCX structure - try different possible paths
                # Handle both forward and backslash separators
                doc_path = None
                for possible_path in ['word/document.xml', 'word\\document.xml', 'Word/document.xml', 'document.xml']:
                    if possible_path in file_list:
                        doc_path = possible_path
                        break
                if not doc_path:
                    # Log more details about what we found
                    self.log_error(f"No document.xml found. ZIP contains: {file_list}")
                    raise ValueError(f"Not a valid DOCX file structure. Files in archive: {len(file_list)}")
                # Extract the main document XML
                with zip_file.open(doc_path) as xml_file:
                    tree = ET.parse(xml_file)
                    root = tree.getroot()
                    # Define namespaces
                    namespaces = {
                        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
                    }
                    # Extract all text from paragraphs
                    content_parts = []
                    for paragraph in root.findall('.//w:p', namespaces):
                        texts = []
                        for text_elem in paragraph.findall('.//w:t', namespaces):
                            if text_elem.text:
                                texts.append(text_elem.text)
                        if texts:
                            para_text = ''.join(texts).strip()
                            if para_text:
                                content_parts.append(para_text)
                    # Try to extract tables
                    for table in root.findall('.//w:tbl', namespaces):
                        table_rows = []
                        for row in table.findall('.//w:tr', namespaces):
                            cells = []
                            for cell in row.findall('.//w:tc', namespaces):
                                cell_texts = []
                                for text_elem in cell.findall('.//w:t', namespaces):
                                    if text_elem.text:
                                        cell_texts.append(text_elem.text)
                                cells.append(' '.join(cell_texts))
                            if cells:
                                table_rows.append(' | '.join(cells))
                        if table_rows:
                            content_parts.append('\n[Table]\n' + '\n'.join(table_rows))
                    content = '\n\n'.join(content_parts)
                    return {
                        'content': content or '[No text content extracted]',
                        'format': 'docx',
                        'extraction_method': 'xml_parsing',
                        'warning': 'Document was parsed using alternative method - formatting may be simplified'
                    }
        except Exception as e:
            self.log_error(f"Alternative DOCX parsing also failed: {e}")
            # Final fallback - try to extract any text using textract or similar
            return self.parse_docx_textract_fallback(file_data, filename)
    def parse_docx_textract_fallback(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """
        Final DOCX fallback: pandoc reading the document from STDIN.
        No disk access - the file bytes are piped straight to the subprocess.
        """
        try:
            result = subprocess.run(
                ['pandoc', '-f', 'docx', '-t', 'plain'],
                input=file_data,
                capture_output=True,
                timeout=10
            )
            if result.returncode == 0 and result.stdout:
                return {
                    'content': result.stdout.decode('utf-8', errors='replace'),
                    'format': 'docx',
                    'extraction_method': 'pandoc',
                    'warning': 'Document extracted using pandoc - formatting simplified'
                }
        except (subprocess.SubprocessError, FileNotFoundError) as e:
            self.log_error(f"Pandoc fallback failed: {e}")
        # Absolute final fallback
        return {
            'content': f'[Unable to extract content from {filename} - document format is not supported or file is corrupted]',
            'format': 'docx',
            'extraction_method': 'failed',
            'error': 'All extraction methods failed'
        }
    def parse_excel_adaptive(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """
        Extract Excel/CSV content with adaptive structure preservation
        Args:
            file_data: Excel/CSV file data as bytes
            filename: Original filename for error reporting
        Returns:
            Dict with sheets, content, and structure info
        """
        try:
            # Determine if this is a delimited-text table or a binary spreadsheet.
            # .tsv/.tab (tab) and .psv (pipe) go through the CSV path — pandas sniffs
            # the delimiter across [',', ';', '\t', '|'] in _parse_csv_with_pandas, so
            # routing them here is deliberate rather than relying on the except fallback.
            if filename.lower().endswith(('.csv', '.tsv', '.tab', '.psv')):
                return self._parse_csv_with_pandas(file_data, filename, pd)
            else:
                return self._parse_excel_with_pandas(file_data, filename, pd)
        except Exception as e:
            # Try CSV interpretation as fallback
            try:
                self.log_warning(f"Excel parsing failed for {filename}, trying CSV fallback: {e}")
                return self._parse_csv_with_pandas(file_data, filename, pd)
            except Exception as csv_error:
                raise Exception(f"Spreadsheet parsing failed: {e}, CSV fallback: {csv_error}")
    def _parse_excel_with_pandas(self, file_data: bytes, filename: str, pd) -> Dict[str, Any]:
        """Parse Excel file using pandas"""
        # Read all sheets
        try:
            all_sheets = pd.read_excel(BytesIO(file_data), sheet_name=None, engine='openpyxl')
        except Exception:
            # Try without specifying engine
            all_sheets = pd.read_excel(BytesIO(file_data), sheet_name=None)
        sheets = []
        content_parts = []
        # Limit number of sheets
        sheet_items = list(all_sheets.items())[:20]
        for sheet_name, df in sheet_items:
            sheet_data = {
                'name': str(sheet_name)[:50],  # Limit sheet name length
                'rows': len(df),
                'cols': len(df.columns)
            }
            # Handle different data structures
            if df.empty:
                sheet_content = "[Empty sheet]"
                sheet_data['format'] = 'empty'
            elif self._is_simple_table(df):
                # Convert to markdown table
                sheet_content = self._dataframe_to_markdown(df)
                sheet_data['format'] = 'table'
            elif self._is_list_data(df):
                # Simple list format
                sheet_content = self._format_as_list(df)
                sheet_data['format'] = 'list'
            else:
                # Fallback: structured text
                sheet_content = df.to_string(max_rows=1000, max_cols=50)
                sheet_data['format'] = 'raw'
            sheet_data['content'] = sheet_content
            # Add to overall content
            content_parts.append(f"[Sheet: {sheet_name}]")
            content_parts.append(sheet_content)
            sheets.append(sheet_data)
        return {
            'content': '\n\n'.join(content_parts),
            'sheets': sheets,
            'total_sheets': len(all_sheets),
            'format': 'excel'
        }
    def _parse_csv_with_pandas(self, file_data: bytes, filename: str, pd) -> Dict[str, Any]:
        """Parse CSV file using pandas"""
        try:
            # Try different encodings and separators
            text_data = file_data.decode('utf-8')
        except UnicodeDecodeError:
            try:
                text_data = file_data.decode('latin-1')
            except UnicodeDecodeError:
                text_data = file_data.decode('utf-8', errors='ignore')
        # Try to read CSV with different separators
        for sep in [',', ';', '\t', '|']:
            try:
                df = pd.read_csv(BytesIO(text_data.encode('utf-8')), 
                               sep=sep, on_bad_lines='skip', nrows=10000)
                if len(df.columns) > 1:  # Found good separator
                    break
            except Exception:
                continue
        else:
            # Fallback: assume comma separation
            df = pd.read_csv(BytesIO(text_data.encode('utf-8')), 
                           on_bad_lines='skip', nrows=10000)
        # Format the data
        if df.empty:
            content = "[Empty CSV file]"
        else:
            content = self._dataframe_to_markdown(df)
        return {
            'content': content,
            'sheets': [{'name': 'CSV Data', 'content': content, 'format': 'table'}],
            'total_sheets': 1,
            'rows': len(df),
            'cols': len(df.columns),
            'format': 'csv'
        }
    def parse_text(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """
        Parse plain text files with encoding detection
        Args:
            file_data: Text file data as bytes
            filename: Original filename for error reporting
        Returns:
            Dict with text content
        """
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        for encoding in encodings:
            try:
                text = file_data.decode(encoding)
                # Detect file type from content or extension
                file_format = self._detect_text_format(text, filename)
                return {
                    'content': text,
                    'format': file_format,
                    'encoding': encoding,
                    'lines': len(text.splitlines())
                }
            except UnicodeDecodeError:
                continue
        # If all encodings fail, use utf-8 with error handling
        text = file_data.decode('utf-8', errors='ignore')
        return {
            'content': text,
            'format': 'text',
            'encoding': 'utf-8_with_errors',
            'lines': len(text.splitlines()),
            'warning': 'Some characters may not have been decoded correctly'
        }
    def parse_rtf(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract prose from an RTF document, stripping control words.

        Without this, .rtf fell through to parse_text and handed the model the raw
        ``{\\rtf1\\ansi ...`` control-code soup. striprtf (BSD-3) consumes a decoded
        string, so the bytes are decoded first (utf-8, then latin-1 as a fallback).
        """
        from striprtf.striprtf import rtf_to_text
        try:
            raw = file_data.decode('utf-8')
        except UnicodeDecodeError:
            raw = file_data.decode('latin-1', errors='replace')
        text = rtf_to_text(raw)
        return {
            'content': text or '[No text content found in RTF document]',
            'format': 'rtf',
            'lines': len((text or '').splitlines()),
        }
    def parse_email(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract an .eml message: key headers plus the text/plain body.

        Stdlib ``email`` only — no dependency. Attachments are intentionally
        skipped: only the human-readable headers (From/To/Cc/Subject/Date) and the
        plain-text body part reach the model.
        """
        from email import policy
        from email.parser import BytesParser
        msg = BytesParser(policy=policy.default).parsebytes(file_data)
        header_lines = []
        for header in ('From', 'To', 'Cc', 'Subject', 'Date'):
            value = msg.get(header)
            if value:
                header_lines.append(f"{header}: {value}")
        body = ''
        try:
            body_part = msg.get_body(preferencelist=('plain',))
            if body_part is not None:
                body = body_part.get_content()
        except Exception as e:
            self.log_warning(f"Failed to extract .eml body from {filename}: {e}")
        content = '\n'.join(header_lines)
        if body:
            content += ('\n\n' if content else '') + body.strip()
        return {
            'content': content or '[No readable headers or text body in email]',
            'format': 'eml',
        }
    def parse_notebook(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract a Jupyter .ipynb: markdown and code cells only.

        Cell OUTPUTS, execution counts, and base64-embedded image payloads are
        deliberately dropped — a notebook's outputs can be megabytes of base64 that
        would blow the context window. This is why .ipynb has a dedicated handler and
        never falls through to raw parse_text.
        """
        import json
        nb = json.loads(file_data.decode('utf-8', errors='replace'))
        cells = nb.get('cells', [])
        blocks = []
        for cell in cells:
            if not isinstance(cell, dict):
                continue
            cell_type = cell.get('cell_type')
            source = cell.get('source', '')
            if isinstance(source, list):
                source = ''.join(source)
            source = (source or '').strip()
            if not source:
                continue
            if cell_type == 'markdown':
                blocks.append(source)
            elif cell_type == 'code':
                blocks.append(f"```\n{source}\n```")
        content = '\n\n'.join(blocks)
        return {
            'content': content or '[Notebook contains no markdown or code cells]',
            'format': 'ipynb',
            'total_cells': len(cells),
        }
    def flexible_table_to_markdown(self, table_data: List[List[str]]) -> str:
        """
        Convert variable table structures to markdown with error handling
        Args:
            table_data: 2D list representing table rows and columns
        Returns:
            Markdown table string or empty string if conversion fails
        """
        if not table_data or not table_data[0]:
            return ""
        try:
            # Handle irregular tables gracefully
            max_cols = max(len(row) if row else 0 for row in table_data)
            if max_cols == 0:
                return ""
            # Normalize rows to same width
            normalized = []
            for row in table_data:
                if row:
                    # Pad short rows, handle None values
                    norm_row = [(str(cell) if cell is not None else "") for cell in row]
                    norm_row.extend([""] * (max_cols - len(norm_row)))
                    normalized.append(norm_row)
            if not normalized:
                return ""
            # Build markdown table
            lines = []
            # Header row
            header = "| " + " | ".join(cell.replace("|", "\\|") for cell in normalized[0]) + " |"
            lines.append(header)
            # Separator
            separator = "|" + "|".join([" --- " for _ in range(max_cols)]) + "|"
            lines.append(separator)
            # Data rows (limit to reasonable number)
            for row in normalized[1:100]:  # Limit to 100 rows
                line = "| " + " | ".join(cell.replace("|", "\\|") for cell in row) + " |"
                lines.append(line)
            if len(normalized) > 101:
                lines.append(f"| ... | ({len(normalized) - 101} more rows) | ... |")
            return "\n".join(lines)
        except Exception as e:
            self.log_warning(f"Table conversion failed: {e}")
            return "[Table data could not be converted]"
    def fix_markdown_tables(self, text: str) -> str:
        """
        Fix common markdown table formatting issues
        Args:
            text: Text containing potentially malformed markdown tables
        Returns:
            Text with fixed table formatting
        """
        try:
            lines = text.split('\n')
            fixed_lines = []
            in_table = False
            table_col_count = 0
            for line in lines:
                if '|' in line and line.strip():
                    # This looks like a table row
                    stripped = line.strip()
                    # Ensure proper table row format
                    if not stripped.startswith('|'):
                        stripped = '| ' + stripped
                    if not stripped.endswith('|'):
                        stripped = stripped + ' |'
                    # Count columns
                    col_count = stripped.count('|') - 1
                    if not in_table:
                        # Starting a new table
                        table_col_count = col_count
                        in_table = True
                    elif col_count != table_col_count:
                        # Column count mismatch - normalize
                        stripped = self._normalize_table_row(stripped, table_col_count)
                    fixed_lines.append(stripped)
                elif in_table and line.strip() and '|' not in line:
                    # Table ended
                    in_table = False
                    fixed_lines.append(line)
                else:
                    # Regular line
                    if in_table and not line.strip():
                        in_table = False  # Empty line ends table
                    fixed_lines.append(line)
            return '\n'.join(fixed_lines)
        except Exception as e:
            self.log_warning(f"Table fixing failed: {e}")
            return text
    def force_text_extraction(self, file_data: bytes, mime_type: str, filename: str) -> str:
        """
        Last resort text extraction for when structured parsing fails
        Args:
            file_data: File data as bytes
            mime_type: MIME type of the file
            filename: Original filename
        Returns:
            Extracted text content
        """
        # Don't try to decode binary formats as text
        binary_extensions = ('.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.pdf', '.zip', '.rar')
        if filename.lower().endswith(binary_extensions):
            # These are binary formats that shouldn't be decoded as text
            return f"[Unable to extract text from corrupted {filename} - binary format]"
        # Try direct text decoding for text-based formats
        encodings = ['utf-8', 'latin-1', 'cp1252']
        for encoding in encodings:
            try:
                text = file_data.decode(encoding, errors='ignore')
                # Clean up the text
                text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
                # Check if we got mostly readable text (not binary garbage)
                if text.strip() and len([c for c in text[:100] if c.isprintable()]) > 80:
                    return f"[Raw text extraction from {filename}]\n{text}"
            except Exception:
                continue
        return f"[Unable to extract readable text from {filename}]"
    # Helper methods
    def _detect_text_structure(self, text: str) -> Dict[str, Any]:
        """Detect structural elements in text"""
        lines = text.split('\n')
        structure = {
            'has_headers': False,
            'has_lists': False,
            'has_tables': False,
            'line_count': len(lines)
        }
        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue
            # Check for headers (short lines, all caps, etc.)
            if len(stripped) < 100 and (stripped.isupper() or 
                                       stripped.startswith('#') or
                                       len(stripped.split()) <= 5):
                structure['has_headers'] = True
            # Check for lists
            if re.match(r'^\s*[-*•]\s+', line) or re.match(r'^\s*\d+\.\s+', line):
                structure['has_lists'] = True
            # Check for table-like structures
            if '|' in line and line.count('|') >= 2:
                structure['has_tables'] = True
        return structure
    def _extract_docx_table(self, table) -> str:
        """Extract table from docx document as markdown"""
        try:
            rows = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    # Clean cell text
                    cell_text = cell.text.strip().replace('\n', ' ')
                    cells.append(cell_text)
                rows.append(cells)
            return self.flexible_table_to_markdown(rows)
        except Exception as e:
            self.log_warning(f"Failed to extract docx table: {e}")
            return "[Table extraction failed]"
    def _split_into_sections(self, content: str) -> List[Dict[str, str]]:
        """Split content into sections based on headers"""
        lines = content.split('\n')
        sections = []
        current_section = {'title': 'Document Start', 'content': []}
        for line in lines:
            if line.startswith('##'):
                # New section
                if current_section['content']:
                    current_section['content'] = '\n'.join(current_section['content']).strip()
                    sections.append(current_section)
                current_section = {
                    'title': line.replace('##', '').strip(),
                    'content': []
                }
            else:
                current_section['content'].append(line)
        # Add final section
        if current_section['content']:
            current_section['content'] = '\n'.join(current_section['content']).strip()
            sections.append(current_section)
        return sections
    def _is_simple_table(self, df) -> bool:
        """Check if DataFrame is a simple table suitable for markdown"""
        if df.empty or len(df.columns) < 2:
            return False
        # Check if column names look reasonable
        unnamed_cols = sum(1 for col in df.columns if 'Unnamed' in str(col))
        if unnamed_cols > len(df.columns) / 2:
            return False
        return True
    def _is_list_data(self, df) -> bool:
        """Check if DataFrame is better represented as a list"""
        return len(df.columns) <= 2 and len(df) > 10
    def _dataframe_to_markdown(self, df) -> str:
        """Convert DataFrame to markdown with size limits"""
        try:
            # Limit size
            if len(df) > 1000:
                df_limited = df.head(1000)
                truncated = True
            else:
                df_limited = df
                truncated = False
            if len(df_limited.columns) > 20:
                df_limited = df_limited.iloc[:, :20]
                cols_truncated = True
            else:
                cols_truncated = False
            # Clean column names
            df_limited.columns = [str(col)[:50] for col in df_limited.columns]
            # Convert to markdown
            md_table = df_limited.to_markdown(index=False, tablefmt='pipe')
            # Add truncation notices
            notices = []
            if truncated:
                notices.append(f"[Showing first 1000 of {len(df)} rows]")
            if cols_truncated:
                notices.append(f"[Showing first 20 of {len(df.columns)} columns]")
            if notices:
                md_table = '\n'.join(notices) + '\n\n' + md_table
            return md_table
        except Exception as e:
            self.log_warning(f"DataFrame to markdown conversion failed: {e}")
            return df.to_string(max_rows=100, max_cols=10)
    def _format_as_list(self, df) -> str:
        """Format DataFrame as a simple list"""
        try:
            lines = []
            for _, row in df.head(1000).iterrows():
                if len(df.columns) == 1:
                    lines.append(f"• {row.iloc[0]}")
                else:
                    lines.append(f"• {row.iloc[0]}: {row.iloc[1]}")
            if len(df) > 1000:
                lines.append(f"[... and {len(df) - 1000} more items]")
            return '\n'.join(lines)
        except Exception as e:
            self.log_warning(f"List formatting failed: {e}")
            return df.to_string()
    def _detect_text_format(self, text: str, filename: str) -> str:
        """Detect the format of a text file"""
        filename_lower = filename.lower()
        if filename_lower.endswith('.py'):
            return 'python'
        elif filename_lower.endswith('.js'):
            return 'javascript'
        elif filename_lower.endswith('.json'):
            return 'json'
        elif filename_lower.endswith('.sql'):
            return 'sql'
        elif filename_lower.endswith(('.yaml', '.yml')):
            return 'yaml'
        elif filename_lower.endswith('.xml'):
            return 'xml'
        elif filename_lower.endswith(('.html', '.htm')):
            return 'html'
        elif filename_lower.endswith('.md'):
            return 'markdown'
        elif filename_lower.endswith('.log'):
            return 'log'
        else:
            return 'text'
    def _normalize_table_row(self, row: str, target_cols: int) -> str:
        """Normalize table row to have target number of columns"""
        try:
            parts = row.split('|')
            # Remove first and last empty parts (from leading/trailing |)
            if parts and not parts[0].strip():
                parts = parts[1:]
            if parts and not parts[-1].strip():
                parts = parts[:-1]
            # Adjust to target column count
            if len(parts) < target_cols:
                parts.extend([''] * (target_cols - len(parts)))
            elif len(parts) > target_cols:
                parts = parts[:target_cols]
            return '| ' + ' | '.join(parts) + ' |'
        except Exception:
            return row
    def convert_pdf_to_images(self, file_data: bytes, max_pages: int = 10) -> List[Dict[str, Any]]:
        """
        Convert PDF pages to images for vision/OCR processing
        Args:
            file_data: PDF file data as bytes
            max_pages: Maximum number of pages to convert (default 10 - OpenAI limit)
        Returns:
            List of dicts with page images as base64 and metadata
        """
        try:
            # Convert PDF pages to PIL images
            # Note: This requires poppler-utils to be installed on the system
            images = convert_from_bytes(file_data, dpi=150, fmt='png')
            converted_pages = []
            pages_to_process = min(len(images), max_pages)
            for i, image in enumerate(images[:pages_to_process]):
                try:
                    # Convert PIL image to base64
                    buffer = BytesIO()
                    image.save(buffer, format='PNG')
                    image_data = buffer.getvalue()
                    base64_data = base64.b64encode(image_data).decode('utf-8')
                    converted_pages.append({
                        'page': i + 1,
                        'base64_data': base64_data,
                        'mimetype': 'image/png',
                        'width': image.width,
                        'height': image.height
                    })
                    self.log_debug(f"Converted PDF page {i+1} to image ({image.width}x{image.height})")
                except Exception as e:
                    self.log_warning(f"Failed to convert page {i+1} to image: {e}")
                    continue
            if len(images) > max_pages:
                self.log_info(f"Converted first {max_pages} of {len(images)} PDF pages to images")
            else:
                self.log_info(f"Converted all {len(images)} PDF pages to images")
            return converted_pages
        except Exception as e:
            self.log_error(f"Failed to convert PDF to images: {e}")
            # Check if it's a poppler issue
            if "poppler" in str(e).lower():
                self.log_error("poppler-utils may not be installed. Install with: apt-get install poppler-utils (Linux) or brew install poppler (Mac)")
            return []
    def _is_image_based_pdf(self, pdf_result: Dict[str, Any]) -> bool:
        """
        Detect if a PDF is likely image-based (scanned document)
        Args:
            pdf_result: Result from PDF parsing
        Returns:
            True if PDF appears to be image-based/scanned
        """
        # Check if we have pages data
        pages = pdf_result.get('pages', [])
        if not pages:
            return False
        # Count pages with meaningful text content
        pages_with_text = 0
        total_text_length = 0
        for page in pages:
            content = page.get('content', '')
            # Remove whitespace and check length
            clean_content = content.strip()
            # Skip page markers and extraction failure messages
            if clean_content and not clean_content.startswith('[') and len(clean_content) > 50:
                pages_with_text += 1
                total_text_length += len(clean_content)
        # Heuristics to determine if PDF is image-based:
        # 1. Less than 20% of pages have meaningful text
        # 2. Average text per page is very low (< 100 chars)
        total_pages = pdf_result.get('total_pages', len(pages))
        if total_pages == 0:
            return False
        text_page_ratio = pages_with_text / total_pages
        avg_text_per_page = total_text_length / total_pages if total_pages > 0 else 0
        # Consider it image-based if:
        # - Very few pages have text (< 20%)
        # - OR average text per page is very low (< 100 chars)
        # - OR the entire content is very short relative to page count
        is_likely_image_based = (
            text_page_ratio < 0.2 or 
            avg_text_per_page < 100 or
            (total_pages > 1 and total_text_length < 200)
        )
        if is_likely_image_based:
            self.log_info(f"PDF appears to be image-based: {pages_with_text}/{total_pages} pages with text, "
                         f"avg {avg_text_per_page:.0f} chars/page")
        return is_likely_image_based