"""D1 document-handler fixes: pptx support, router unification, timeout,
zip-bomb guard, no-disk rule, and the generated supported-types message."""
import asyncio
import io
import zipfile

import pytest

import document_handler as dh
from document_handler import (
    DOCUMENT_EXTENSIONS,
    MIME_TYPE_HANDLERS,
    SUPPORTED_DOCUMENT_MIMETYPES,
    DocumentHandler,
)

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@pytest.fixture
def handler():
    return DocumentHandler()


def _make_pptx(with_table=True, with_notes=True) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = "Quarterly Review"
    slide.placeholders[1].text = "Revenue grew 12%"
    if with_table:
        rows, cols = 2, 2
        table = slide.shapes.add_table(rows, cols, Inches(1), Inches(3), Inches(4), Inches(1)).table
        table.cell(0, 0).text = "Region"
        table.cell(0, 1).text = "Sales"
        table.cell(1, 0).text = "West"
        table.cell(1, 1).text = "42"
    if with_notes:
        slide.notes_slide.notes_text_frame.text = "Mention the churn caveat"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- pptx support ---

def test_pptx_extraction(handler):
    result = handler.safe_extract_content(_make_pptx(), PPTX_MIME, "review.pptx")
    content = result["content"]
    assert result.get("error") is None or "error" not in result
    assert result["format"] == "pptx"
    assert "Quarterly Review" in content
    assert "Revenue grew 12%" in content
    assert "West" in content and "42" in content            # table cells
    assert "Mention the churn caveat" in content            # speaker notes
    assert "Slide 1" in content


def test_pptx_via_filename_fallback(handler):
    # Unknown mimetype but .pptx filename must still route to the pptx parser
    result = handler.safe_extract_content(_make_pptx(), "application/octet-stream", "deck.pptx")
    assert "Quarterly Review" in result["content"]


def test_legacy_ppt_and_doc_now_unsupported(handler):
    assert ".ppt" not in DOCUMENT_EXTENSIONS
    assert ".doc" not in DOCUMENT_EXTENSIONS
    assert "application/vnd.ms-powerpoint" not in SUPPORTED_DOCUMENT_MIMETYPES
    assert "application/msword" not in SUPPORTED_DOCUMENT_MIMETYPES
    assert not handler.is_document_file("old.ppt", "application/vnd.ms-powerpoint")
    assert not handler.is_document_file("old.doc", "application/msword")


# --- router parity by construction ---

def test_every_supported_mimetype_has_a_real_handler():
    """SUPPORTED_DOCUMENT_MIMETYPES ⊆ MIME_TYPE_HANDLERS, and every handler exists.
    This is what prevented the old async/sync router drift from being caught."""
    h = DocumentHandler()
    missing = SUPPORTED_DOCUMENT_MIMETYPES - set(MIME_TYPE_HANDLERS)
    assert not missing, f"Supported mimetypes with no routing entry: {missing}"
    for mimetype, handler_name in MIME_TYPE_HANDLERS.items():
        method = getattr(h, handler_name, None)
        assert callable(method), f"{mimetype} routes to nonexistent handler {handler_name}"


def test_async_is_a_wrapper_of_sync_router():
    """The async entry must delegate to the sync implementation (single code path)."""
    import inspect
    src = inspect.getsource(DocumentHandler.safe_extract_content_async)
    assert "self.safe_extract_content" in src
    # No independent routing logic allowed in the async path
    assert "mime_type ==" not in src and "endswith" not in src


@pytest.mark.asyncio
async def test_async_returns_sync_metadata_and_sanitization(handler):
    # html was one of the types the drifted async router dead-ended
    result = await handler.safe_extract_content_async(
        b"<html><body>hello</body></html>", "text/html", "page.html"
    )
    assert result["filename"] == "page.html"
    assert result["mime_type"] == "text/html"
    assert result["size_bytes"] > 0
    assert "hello" in result["content"]


@pytest.mark.asyncio
async def test_csv_extraction_works_via_async(handler):
    # The old async router routed CSV to a phantom method (AttributeError)
    result = await handler.safe_extract_content_async(
        b"name,qty\nwidget,7\n", "text/csv", "inv.csv"
    )
    assert "widget" in result["content"]
    assert result.get("format") != "error"


# --- timeout ---

@pytest.mark.asyncio
async def test_extraction_timeout_returns_placeholder(handler, monkeypatch):
    import time as _time

    def slow_extract(*a, **k):
        _time.sleep(1.5)
        return {"content": "late"}

    monkeypatch.setattr(handler, "safe_extract_content", slow_extract)
    monkeypatch.setattr(dh, "EXTRACTION_TIMEOUT_SECONDS", 0.2)
    result = await handler.safe_extract_content_async(b"x", "application/pdf", "slow.pdf")
    assert result["format"] == "error"
    assert "timed out" in result["content"]
    assert result["filename"] == "slow.pdf"


# --- zip-bomb guard ---

def _make_zip_bomb() -> bytes:
    """Small archive that declares a huge decompressed size."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("word/document.xml", b"\x00" * (250 * 1024 * 1024))
    return buf.getvalue()


def test_zip_bomb_refused(handler):
    bomb = _make_zip_bomb()
    assert len(bomb) < 5 * 1024 * 1024  # it really is small on the wire
    result = handler.safe_extract_content(
        bomb, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "bomb.docx"
    )
    assert result["format"] == "error"
    assert "safe limits" in result["content"]


def test_normal_office_zip_passes_guard(handler):
    assert handler._office_zip_within_limits(_make_pptx()) is True


# --- no-disk rule ---

def test_no_tempfile_usage_in_document_handler():
    """Hard rule: extraction never touches disk. No allowlisted exceptions in
    this module (pdf2image's internal poppler temp files live outside it and
    that path is slated for retirement)."""
    import inspect
    src = inspect.getsource(dh)
    for needle in ("tempfile", "NamedTemporaryFile", "mkstemp", "mkdtemp"):
        # comments/docstrings may mention the rule; strip lines that are comments
        code_lines = [
            line for line in src.splitlines()
            if needle in line and not line.lstrip().startswith("#")
            and '"' not in line and "'" not in line or f"import {needle}" in line
        ]
        assert not [l for l in code_lines if "import" in l], f"disk-touching usage: {code_lines}"


def test_pandoc_fallback_uses_stdin():
    import inspect
    src = inspect.getsource(DocumentHandler.parse_docx_textract_fallback)
    assert "input=file_data" in src
    assert "tempfile" not in src


# --- generated supported-types message ---

def test_unsupported_message_generated_from_extensions():
    import inspect
    from message_processor import base as mp_base
    src = inspect.getsource(mp_base)
    assert "DOCUMENT_EXTENSIONS" in src, "supported-types message must be generated, not hardcoded"
    assert "PDF, DOCX, XLSX, CSV, TXT, etc." not in src


def test_generated_doc_list_matches_sets():
    doc_types = ", ".join(sorted(ext.lstrip(".").upper() for ext in DOCUMENT_EXTENSIONS))
    assert "PPTX" in doc_types and "PDF" in doc_types
    assert "PPT," not in doc_types and not doc_types.endswith("PPT")
    assert "DOC," not in doc_types  # only DOCX remains
